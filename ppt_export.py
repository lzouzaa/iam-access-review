"""
Generate executive PowerPoint presentation from review results.
"""
from io import BytesIO
from typing import Dict, List

from pptx import Presentation
from pptx.util import Inches, Pt

from comparator import (
    CLASS_ACESSO_OBSOLETO,
    CLASS_OK,
    CLASS_POSSIVEL_REMOCAO,
    CLASS_REMOVER_PRIORIDADE,
    CLASS_REVISAO_MANUAL,
)
from i18n import classification_label


def build_executive_pptx(
    stats,
    summary_lines: List[str],
    lang: str,
    date_str: str,
    remove_by_application: Dict[str, int] | None = None,
) -> bytes:
    """
    Build a .pptx with title slide, metrics slide, and summary slide.
    Returns file content as bytes.
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    title_slide = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide)
    title = "Revisão de Acessos IAM" if lang == "pt-BR" else "IAM Access Review"
    subtitle = date_str
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

    # Slide 2: Metrics (table-like as bullets)
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    metrics_title = "Métricas por classificação" if lang == "pt-BR" else "Metrics by classification"
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = metrics_title
    p.font.size = Pt(24)
    p.font.bold = True

    left = Inches(0.5)
    top = Inches(1.4)
    width = Inches(9)
    height = Inches(2.5)
    body = slide.shapes.add_textbox(left, top, width, height)
    tf = body.text_frame
    tf.word_wrap = True
    for cl in [CLASS_OK, CLASS_REMOVER_PRIORIDADE, CLASS_POSSIVEL_REMOCAO, CLASS_REVISAO_MANUAL, CLASS_ACESSO_OBSOLETO]:
        count = int(stats.get(cl, 0))
        label = classification_label(cl, lang)
        p = tf.add_paragraph()
        p.text = f"  • {label}: {count}"
        p.font.size = Pt(14)
        p.space_after = Pt(6)

    # Slide 3: Executive summary / recommendations
    slide = prs.slides.add_slide(layout)
    sum_title = "Resumo executivo e recomendações" if lang == "pt-BR" else "Executive summary and recommendations"
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = sum_title
    p.font.size = Pt(24)
    p.font.bold = True

    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for line in summary_lines:
        if line.strip():
            p = tf.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(12)
            p.space_after = Pt(4)

    # Slide 4: Tools/systems with users recommended for removal
    slide = prs.slides.add_slide(layout)
    rem_title = (
        "Ferramentas com usuários recomendados para remoção"
        if lang == "pt-BR"
        else "Tools with users recommended for removal"
    )
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = rem_title
    p.font.size = Pt(24)
    p.font.bold = True

    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True

    items = []
    if remove_by_application:
        items = sorted(remove_by_application.items(), key=lambda kv: (-kv[1], kv[0]))

    if not items:
        text = (
            "Nenhum usuário recomendado para remoção prioritária."
            if lang == "pt-BR"
            else "No users recommended for priority removal."
        )
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
    else:
        intro = (
            "Quantidade de usuários classificados como 'Remover com prioridade' por ferramenta:"
            if lang == "pt-BR"
            else "Number of users classified as 'Remove with priority' per tool:"
        )
        p = tf.paragraphs[0]
        p.text = intro
        p.font.size = Pt(13)
        p.space_after = Pt(8)

        for app_name, count in items:
            p = tf.add_paragraph()
            p.text = f"  • {app_name}: {int(count)}"
            p.font.size = Pt(13)
            p.space_after = Pt(5)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
